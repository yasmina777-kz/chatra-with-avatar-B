
import io
import logging
import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from uuid import uuid4

logger = logging.getLogger(__name__)

UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")
APP_BASE_URL = os.getenv("APP_BASE_URL", "http://localhost:8000")


class SlideData:
    def __init__(self, index: int, text: str, image_url: str | None = None):
        self.index = index
        self.text = text
        self.image_url = image_url

    def to_dict(self) -> dict:
        return {"index": self.index, "text": self.text, "image_url": self.image_url}


def extract_slides(data: bytes, filename: str) -> list[SlideData]:
    """Главная точка входа. Возвращает список слайдов с текстом и (если получилось) картинкой."""
    ext = Path(filename).suffix.lower()

    if ext == ".pptx":
        texts = _extract_pptx_text(data)
    elif ext == ".pdf":
        texts = _extract_pdf_text(data)
    else:
        raise ValueError(f"Для лекций поддерживаются только .pptx и .pdf, получено: {ext}")

    images = []
    try:
        images = _render_slide_images(data, filename, count=len(texts))
    except Exception as exc:
        logger.warning("Не удалось отрендерить картинки слайдов (%s), продолжаем без них: %s", filename, exc)

    slides = []
    for i, text in enumerate(texts):
        img = images[i] if i < len(images) else None
        slides.append(SlideData(index=i, text=text, image_url=img))

    return slides


def _extract_pptx_text(data: bytes) -> list[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []

    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        parts.append(row_text)

        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Заметки: {notes}]")
        except Exception:
            pass

        texts.append("\n".join(parts).strip() or "(слайд без текста)")

    if not texts:
        raise ValueError("В презентации не найдено слайдов")

    return texts


def _extract_pdf_text(data: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    texts: list[str] = []

    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if not text:
            text = _ocr_pdf_page(page)
        texts.append(text or "(страница без текста)")

    if not texts:
        raise ValueError("В PDF не найдено страниц")

    return texts


def _ocr_pdf_page(page) -> str:
    try:
        import pytesseract
        pil_image = page.to_image(resolution=150).original
        return pytesseract.image_to_string(pil_image, lang="rus+eng").strip()
    except Exception as exc:
        logger.debug("OCR страницы PDF не удался: %s", exc)
        return ""


def _soffice_available() -> bool:
    return shutil.which("soffice") is not None or shutil.which("libreoffice") is not None


def _resolve_executable(name: str) -> str:

    found = shutil.which(name) or shutil.which(name + ".exe")
    return found or name


def _render_slide_images(data: bytes, filename: str, count: int) -> list[str]:

    if not _soffice_available():
        logger.warning("LibreOffice (soffice) не найден в PATH — рендер слайдов в картинки пропущен")
        return []

    soffice_bin = _resolve_executable("soffice")
    pdftoppm_bin = _resolve_executable("pdftoppm")
    logger.info("Рендер слайдов: soffice=%s, pdftoppm=%s, filename=%s", soffice_bin, pdftoppm_bin, filename)

    ext = Path(filename).suffix.lower()

    with tempfile.TemporaryDirectory() as tmpdir:
        src_path = os.path.join(tmpdir, f"source{ext}")
        with open(src_path, "wb") as f:
            f.write(data)

        pdf_path = src_path
        if ext != ".pdf":
            try:
                result = subprocess.run(
                    [soffice_bin, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, src_path],
                    check=True, timeout=180, capture_output=True, text=True,
                )
                logger.info("soffice stdout: %s", result.stdout[:500])
            except subprocess.CalledProcessError as exc:
                logger.error(
                    "soffice завершился с ошибкой (code=%s): stdout=%s stderr=%s",
                    exc.returncode, exc.stdout, exc.stderr,
                )
                return []
            except FileNotFoundError as exc:
                logger.error("Не удалось запустить soffice (%s): %s", soffice_bin, exc)
                return []
            except subprocess.TimeoutExpired:
                logger.error("soffice превысил таймаут конвертации в PDF")
                return []


            candidates = list(Path(tmpdir).glob("*.pdf"))
            if not candidates:
                logger.warning("LibreOffice не создал PDF для %s (файлы в tmpdir: %s)",
                                filename, list(Path(tmpdir).iterdir()))
                return []
            pdf_path = str(candidates[0])

        out_prefix = os.path.join(tmpdir, "slide")

        try:
            result = subprocess.run(
                [pdftoppm_bin, "-png", "-r", "110", pdf_path, out_prefix],
                check=True, timeout=180, capture_output=True, text=True,
            )
        except subprocess.CalledProcessError as exc:
            logger.error(
                "pdftoppm завершился с ошибкой (code=%s): stdout=%s stderr=%s",
                exc.returncode, exc.stdout, exc.stderr,
            )
            return []
        except FileNotFoundError as exc:
            logger.error("Не удалось запустить pdftoppm (%s): %s", pdftoppm_bin, exc)
            return []
        except subprocess.TimeoutExpired:
            logger.error("pdftoppm превысил таймаут конвертации в PNG")
            return []

        generated = sorted(Path(tmpdir).glob("slide-*.png"))
        if not generated:
            logger.warning("pdftoppm не создал ни одной картинки (файлы в tmpdir: %s)",
                            list(Path(tmpdir).iterdir()))
            return []

        logger.info("Успешно отрендерено %d картинок слайдов", len(generated))

        os.makedirs(UPLOAD_DIR, exist_ok=True)
        urls: list[str] = []
        batch_id = uuid4().hex[:10]
        for i, src in enumerate(generated):
            dest_name = f"slide_{batch_id}_{i}.png"
            dest_path = os.path.join(UPLOAD_DIR, dest_name)
            shutil.copyfile(src, dest_path)
            urls.append(f"{APP_BASE_URL.rstrip('/')}/uploads/{dest_name}")

        return urls